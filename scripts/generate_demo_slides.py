"""Generate MITRE demo slide deck (PPTX) for Google Slides import."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT_PATHS = [
    Path(__file__).resolve().parents[1] / "docs" / "MITRE_AI_Decision_Pipeline_Demo.pptx",
    Path.home() / "Downloads" / "MITRE_AI_Decision_Pipeline_Demo_GoogleSlides.pptx",
]

FOOTER = "AI-Enabled Decision Pipeline | MITRE Round 2 Demo"
ACCENT = RGBColor(0, 82, 147)
BODY = RGBColor(40, 40, 40)
MUTED = RGBColor(90, 90, 90)


def _set_run(run, size=18, bold=False, color=BODY):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_footer(slide, text=FOOTER):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9.0), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_run(run, size=10, color=MUTED)
    p.alignment = PP_ALIGN.LEFT


def add_title_slide(prs, title, subtitle, tagline=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(1.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=36, bold=True, color=ACCENT)

    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(8.6), Inches(1.2))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = subtitle
    _set_run(srun, size=18, color=BODY)

    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.0), Inches(8.6), Inches(0.8))
        ttf = tag_box.text_frame
        tp = ttf.paragraphs[0]
        trun = tp.add_run()
        trun.text = tagline
        _set_run(trun, size=14, color=MUTED)

    chips = "Input  →  Decision Engine  →  LLM Explanation  →  Metrics"
    chip_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.0), Inches(8.6), Inches(0.5))
    ctf = chip_box.text_frame
    cp = ctf.paragraphs[0]
    crun = cp.add_run()
    crun.text = chips
    _set_run(crun, size=13, bold=True, color=ACCENT)

    add_footer(slide, "Round 2 project presentation")


def add_section_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=28, bold=True, color=ACCENT)

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(8.2), Inches(5.0))
    btf = body_box.text_frame
    btf.word_wrap = True

    for i, bullet in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.space_after = Pt(10)
        para.font.size = Pt(17)
        para.font.color.rgb = BODY

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(8.2), Inches(0.6))
        ntf = note_box.text_frame
        np = ntf.paragraphs[0]
        nrun = np.add_run()
        nrun.text = note
        _set_run(nrun, size=12, bold=True, color=ACCENT)

    add_footer(slide)


def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(8.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=28, bold=True, color=ACCENT)

    for col, (heading, items, x) in enumerate(
        [(left_title, left_items, 0.7), (right_title, right_items, 5.0)]
    ):
        hbox = slide.shapes.add_textbox(Inches(x), Inches(1.35), Inches(4.1), Inches(0.45))
        htf = hbox.text_frame
        hp = htf.paragraphs[0]
        hrun = hp.add_run()
        hrun.text = heading
        _set_run(hrun, size=16, bold=True, color=ACCENT)

        bbox = slide.shapes.add_textbox(Inches(x), Inches(1.85), Inches(4.1), Inches(4.5))
        btf = bbox.text_frame
        btf.word_wrap = True
        for i, item in enumerate(items):
            para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
            para.text = item
            para.level = 0
            para.space_after = Pt(8)
            para.font.size = Pt(15)
            para.font.color.rgb = BODY

    if note:
        note_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.2), Inches(8.6), Inches(0.6))
        ntf = note_box.text_frame
        np = ntf.paragraphs[0]
        nrun = np.add_run()
        nrun.text = note
        _set_run(nrun, size=12, bold=True, color=ACCENT)

    add_footer(slide)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "AI-Enabled Decision Pipeline",
        "Applied AI systems engineering: deterministic decisioning, constrained multi-step LLM explanations, observability, and cloud modernization.",
        "Responsible AI  •  DevSecOps  •  Observability",
    )

    add_two_column_slide(
        prs,
        "Problem & design thesis",
        "The risk",
        [
            "A GPT wrapper looks impressive but is weak for mission systems.",
            "Unpredictable outputs, poor auditability, unclear decision ownership.",
        ],
        "The design thesis",
        [
            "Spring Boot decides. The LLM only explains and recommends.",
            "Decision path stays deterministic, testable, and auditable.",
            "Users get transparency; engineers keep reliability and control.",
        ],
        note="Core principle: AI augments the system — it does not replace systems engineering discipline.",
    )

    add_section_slide(
        prs,
        "Architecture: separated decisioning, AI, and telemetry",
        [
            "React Frontend → Spring Boot Decision Engine → FastAPI AI Service → OpenAI",
            "DecisionEngine produces DecisionContext before any LLM call.",
            "FastAPI runs two sequential steps: explanation, then recommendations.",
            "Prometheus scrapes FastAPI /metrics; Grafana visualizes health and latency.",
            "Docker Compose runs the full local stack for repeatable demos.",
            "Future cloud path: Azure Container Apps with the same service boundaries.",
        ],
    )

    add_section_slide(
        prs,
        "Live demo flow (3 minutes)",
        [
            "1. Submit case data — user enters income, credit score, employment in the UI.",
            "2. Compute decision — Spring Boot applies scoring and hard-stop rules first.",
            "3. Generate AI output — FastAPI explains the decision and returns 3 recommendations.",
            "4. Show resilience — hard-stop REJECT skips the LLM; fallback still returns a decision.",
            "5. Verify telemetry — Prometheus/Grafana show requests, latency, status, and errors.",
        ],
        note="Demo discipline: show the working path first, then explain boundaries, fallbacks, and future work.",
    )

    add_section_slide(
        prs,
        "Why this is more than a GPT wrapper",
        [
            "Deterministic decision core — APPROVE / REVIEW / REJECT from weighted rules, not the model.",
            "Structured DecisionContext — reasonCodes, ruleFactors, nextStepCategory feed the AI layer.",
            "Hard-stop bypass — critical rejections skip FastAPI entirely (deterministic-hard-stop).",
            "Multi-step constrained pipeline — separate explanation and recommendation prompts with validation.",
            "Forbidden response fields — LLM cannot echo or override decision, score, or reason codes.",
            "Layered fallbacks — invalid JSON, API failure, or missing key still returns the same decision.",
        ],
    )

    add_two_column_slide(
        prs,
        "DevSecOps & observability",
        "Engineering habits demonstrated",
        [
            "Docker Compose for repeatable local startup.",
            "Health checks and demo scripts reduce live-demo risk.",
            "Structured logging with request IDs on the AI service.",
            "Azure Terraform skeleton shows a cloud migration runway.",
        ],
        "Metrics that defend the demo",
        [
            "Target health — is FastAPI reachable?",
            "Request volume — are /analyze calls being received?",
            "Status / errors — are failures visible, not hidden?",
            "Latency — is the user path responsive under load?",
        ],
        note="Proof point: metrics move when ./scripts/demo-traffic.sh runs — not just when Grafana opens.",
    )

    add_section_slide(
        prs,
        "Demo scenarios to show the panel",
        [
            "Scenario A — REVIEW path: 65000 income, 705 credit, EMPLOYED → decision before LLM; source: llm.",
            "Scenario B — Hard-stop REJECT: 45000 income, 580 credit → no FastAPI call; source: deterministic-hard-stop.",
            "Scenario C — Invalid input: missing employment status → HTTP 400 from Spring validation.",
            "Scenario D — Fallback: empty OpenAI key or LLM_FORCE_FALLBACK=true → same decision, rule-based text.",
        ],
        note="Lead with Scenario B to prove decision authority, then Scenario A for the full AI path.",
    )

    add_section_slide(
        prs,
        "What works now + where I would take it next",
        [
            "What works: end-to-end local demo, deterministic engine, two-step LLM pipeline, metrics, Grafana.",
            "Limitations I own: FastAPI-only metrics, no production auth, Azure deploy is future work.",
            "Next phase: confidence scoring, prompt regression tests, human review queue, Azure deployment.",
            "Takeaway: I build AI-enabled systems with boundaries, reliability, governance, and user trust.",
        ],
    )

    return prs


def main():
    prs = build_deck()
    for path in OUTPUT_PATHS:
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(path)
        print(f"Saved: {path}")


if __name__ == "__main__":
    main()
